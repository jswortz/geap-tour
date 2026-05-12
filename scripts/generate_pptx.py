"""Generate PowerPoint slide deck from GEAP workshop content."""

import os
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = os.path.join(os.path.dirname(__file__), "..", "docs")
SCREENSHOTS = os.path.join(DOCS, "screenshots")
DIAGRAMS = os.path.join(os.path.dirname(__file__), "..", "diagrams", "outputs")
OUTPUT = os.path.join(DOCS, "slides.pptx")

BLUE = RGBColor(0x1A, 0x73, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
RED = RGBColor(0xD9, 0x30, 0x25)
YELLOW = RGBColor(0xF9, 0xAB, 0x00)
DARK = RGBColor(0x20, 0x21, 0x24)
GRAY = RGBColor(0x5F, 0x63, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
DARK_BLUE = RGBColor(0x17, 0x4E, 0xA6)

FONT = "Calibri"
CODE_FONT = "Consolas"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_bg_gradient(slide, color1, color2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, font_size=18, bold=False, color=DARK,
             font_name=FONT, space_before=Pt(4), bullet=False, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = space_before
    p.level = level
    if bullet:
        p.level = level
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=DARK, title=None, title_color=BLUE, title_size=22):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = FONT
        p.space_after = Pt(8)

    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_before = Pt(6)
        p.level = 0
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(0xDA, 0xDC, 0xE0)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)

    lines = textwrap.dedent(code_text).strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return shape


def add_card(slide, left, top, width, height, title, body, accent_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = FONT

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.font.name = FONT
    p2.space_before = Pt(6)
    return shape


def add_table(slide, left, top, width, height, headers, rows, font_size=14):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = DARK
                p.font.name = FONT
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return table_shape


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


REPO_URL = "https://github.com/jswortz/geap-tour"
GCP_PROJECT = "wortz-project-352116"


def add_repo_link(slide, code_path=None, y_pos=Inches(6.9)):
    if code_path:
        display_text = f"{REPO_URL}/{code_path}"
        prefix = "/tree/main/" if code_path.endswith("/") else "/blob/main/"
        url = f"{REPO_URL}{prefix}{code_path}"
    else:
        display_text = REPO_URL
        url = REPO_URL
    txBox = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(7), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = display_text
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.name = FONT
    run.hyperlink.address = url
    return txBox


def add_console_link(slide, console_url, y_pos=Inches(6.9)):
    txBox = slide.shapes.add_textbox(Inches(7), y_pos, Inches(4), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Open in Cloud Console"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    run.font.name = FONT
    run.hyperlink.address = console_url
    return txBox


def add_logo(slide, color=GRAY):
    add_text(slide, Inches(11), Inches(6.9), Inches(2), Inches(0.4),
             "Google Cloud", font_size=14, color=color, align=PP_ALIGN.RIGHT)


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ===== SLIDE 1: TITLE =====
    s = prs.slides.add_slide(blank)
    set_bg_gradient(s, DARK_BLUE, BLUE)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Google Cloud Workshop", 16, color=YELLOW, bold=True)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(2),
             "Gemini Enterprise\nAgent Platform Tour", 48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.8),
             "Build, Deploy, Govern, Evaluate, and Optimize AI Agents", 24, color=RGBColor(0xCC, 0xCC, 0xCC))
    add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             "ADK Agents  •  MCP Servers  •  Agent Gateway  •  Agent Registry  •  Model Armor  •  Evaluation Pipeline",
             18, color=RGBColor(0x99, 0x99, 0x99))
    add_repo_link(s, y_pos=Inches(6.0))
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "Welcome everyone to the Gemini Enterprise Agent Platform workshop. "
              "Today we'll go from zero to a fully deployed, governed, and evaluated multi-agent system on Google Cloud. "
              "We'll build real ADK agents, connect them via MCP servers, deploy to Agent Runtime, "
              "set up dual-mode gateway governance, run evaluations, and configure Model Armor for content safety. "
              "This is a hands-on workshop — you'll be writing and deploying code throughout.")

    # ===== SLIDE 2: TABLE OF CONTENTS =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Table of Contents", 42, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
             "Four sessions — from building agents to securing them in production", 20, color=GRAY)

    toc_items = [
        (BLUE,   "1", "AI Gateway / MCP Gateway",
         "Multi-agent architecture  •  MCP tool servers  •  ADK agents  •  Deployment  •  Workload identity"),
        (BLUE,   "2", "AI Gateway — Governance & Eval",
         "Agent Gateway (dual mode)  •  Three-tier evaluation  •  Observability  •  GEPA optimization"),
        (YELLOW, "3", "Agent Registry",
         "Agent registration & discovery  •  Toolspec association  •  Lifecycle governance"),
        (RED,    "4", "Model Security / Model Armor",
         "Input/output screening  •  Prompt injection detection  •  Guardrail callbacks"),
    ]
    for i, (color, num, title, desc) in enumerate(toc_items):
        y = Inches(1.8 + i * 1.35)
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Title
        add_text(s, Inches(1.8), y, Inches(10), Inches(0.5),
                 title, 24, bold=True, color=DARK)
        # Description
        add_text(s, Inches(1.8), Emu(y + Inches(0.45)),
                 Inches(10), Inches(0.5),
                 desc, 16, color=GRAY)
    add_logo(s)
    add_notes(s, "Four sessions today. Session 1 is the longest — we'll build the full agent stack from scratch. "
              "Session 2 continues with governance, evaluation, and observability. "
              "Sessions 3 and 4 are shorter focused deep-dives on Agent Registry and Model Armor. "
              "There's a lunch break between sessions 1 and 2, and a short break before the final wrap-up.")

    # ===== SLIDE 3: ARCHITECTURE =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8), "Platform Architecture", 42, bold=True)
    add_image_safe(s, os.path.join(SCREENSHOTS, "geap_architecture.png"),
                   Inches(1.5), Inches(1.4), width=Inches(10))
    add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "User → Frontend → Agent Gateway → Agent Identity (Runtime) → Agent Gateway → Agents, Tools, Models, APIs",
             16, color=GRAY, align=PP_ALIGN.CENTER)
    add_logo(s)
    add_notes(s, "This is the high-level architecture of GEAP. Key point: the Agent Gateway sits at the network boundary — "
              "all traffic between users, agents, and external services flows through it. "
              "The Agent Identity layer (powered by SPIFFE) gives each agent its own cryptographic identity. "
              "This means we can write policies about WHICH agent can call WHICH tool — not just which user.")

    # ===== SLIDE 4: IDENTITY MODEL =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8), "Agent Identity Model", 42, bold=True)
    add_image_safe(s, os.path.join(SCREENSHOTS, "identity_types.png"),
                   Inches(1), Inches(1.3), width=Inches(11))
    add_logo(s)
    add_notes(s, "Three identity types are critical to understand. "
              "ID-1 is the human user's identity — standard OAuth/OIDC. "
              "ID-2 is the agent's own identity — a SPIFFE ID issued at deploy time. "
              "ID-3 is delegated — the agent acting on behalf of a user with constrained permissions. "
              "This tri-modal identity model is what makes fine-grained governance possible.")

    # ===== SLIDE 5: SESSION 1 HEADER =====
    s = prs.slides.add_slide(blank)
    set_bg(s, BLUE)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Session 1", 16, color=RGBColor(0xCC, 0xDD, 0xFF), bold=True)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "AI Gateway / MCP Gateway", 48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.6),
             "Build, connect, and deploy agents with MCP tool servers", 22, color=RGBColor(0xCC, 0xDD, 0xFF))
    add_bullet_list(s, Inches(0.8), Inches(4.0), Inches(10), Inches(3),
                    ["Multi-agent architecture with MCP connectivity",
                     "FastMCP tool servers on Cloud Run",
                     "ADK agents deployed to Agent Runtime",
                     "SPIFFE-based workload identity"],
                    font_size=20, color=RGBColor(0xEE, 0xEE, 0xFF))
    add_repo_link(s, y_pos=Inches(6.5))
    add_logo(s, RGBColor(0x88, 0x99, 0xCC))
    add_notes(s, "Session 1 is where we build everything. By the end of this session you'll have: "
              "three MCP tool servers running on Cloud Run, three ADK agents deployed to Agent Runtime, "
              "and workload identity configured for secure agent-to-service communication. "
              "We'll start with the architecture, then write code, then deploy.")

    # ===== SLIDE 6: MULTI-AGENT TOPOLOGY =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Multi-Agent Topology", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "01_multi_agent_topology.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_bullet_list(s, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5),
                    ["Coordinator Agent — routes requests to specialists",
                     "Travel Agent — searches flights/hotels, makes bookings",
                     "Expense Agent — submits expenses, enforces policy limits"],
                    title="Three ADK Agents", title_color=BLUE)
    add_bullet_list(s, Inches(7), Inches(4.2), Inches(5.5), Inches(2.5),
                    ["Search MCP — shared by Travel + Coordinator",
                     "Booking MCP — flight/hotel reservations",
                     "Expense MCP — expense submission + policy checks"],
                    title="Three MCP Servers", title_color=BLUE)
    add_logo(s)
    add_notes(s, "Our workshop builds a corporate travel assistant. The Coordinator agent routes requests "
              "to either the Travel agent (flights, hotels) or the Expense agent (reimbursements). "
              "Notice the MCP servers are shared — the Search MCP is used by both Travel and Coordinator. "
              "This is a key MCP advantage: tools are deployed once and shared across agents.")

    # ===== SLIDE 7: MCP SERVER DEV =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "MCP Server Development", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4),
             "FastMCP Pattern", 22, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.2),
                   'from fastmcp import FastMCP\n\nmcp = FastMCP("search-mcp")\n\n@mcp.tool()\ndef search_flights(\n    origin: str,\n    destination: str,\n    date: str | None = None\n) -> list[dict]:\n    """Search available flights."""\n    return flight_db.search(\n        origin, destination, date\n    )',
                   font_size=12)
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(1.4),
             "StreamableHTTP Transport",
             "MCP servers run as standard HTTP services on Cloud Run with automatic scaling", BLUE)
    add_card(s, Inches(7), Inches(3.0), Inches(5.5), Inches(1.4),
             "Tool Discovery",
             "Agents dynamically discover available tools via MCP protocol introspection", GREEN)
    add_card(s, Inches(7), Inches(4.7), Inches(5.5), Inches(1.4),
             "1-to-Many Topology",
             "Multiple agents can share the same MCP server — no duplication needed", YELLOW)
    add_repo_link(s, "src/mcp_servers/search/server.py")
    add_logo(s)
    add_notes(s, "FastMCP makes building tool servers trivial — it's a single decorator on a Python function. "
              "The @mcp.tool() decorator handles schema generation, input validation, and protocol compliance. "
              "We use StreamableHTTP transport so these run as standard HTTP services on Cloud Run. "
              "DEMO: Walk through the search_flights tool code and show how MCP introspection exposes the tool schema.")

    # ===== SLIDE 8: ADK AGENTS =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "ADK Agent Architecture", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4),
             "Agent Definition", 22, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
                   'from google.adk.agents import LlmAgent\nfrom google.adk.integrations\\\n  .agent_registry import AgentRegistry\n\nregistry = AgentRegistry(\n  project_id=PROJECT, location="global"\n)\n\ntravel_agent = LlmAgent(\n  name="travel_agent",\n  model="gemini-2.0-flash",\n  instruction="""You are a travel\n  assistant...""",\n  tools=[\n    registry.get_mcp_toolset(\n      SEARCH_MCP_SERVER\n    )\n  ]\n)',
                   font_size=11)
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(1.4),
             "Coordinator Pattern",
             "Root agent uses sub_agents=[travel, expense] to delegate by intent", BLUE)
    add_card(s, Inches(7), Inches(3.0), Inches(5.5), Inches(1.4),
             "Session Management",
             "Each user gets isolated sessions with conversation history and state", GREEN)
    add_card(s, Inches(7), Inches(4.7), Inches(5.5), Inches(1.4),
             "OTel Tracing",
             "Every agent call produces OpenTelemetry spans — exported to Cloud Trace automatically", YELLOW)
    add_repo_link(s, "src/agents/coordinator_agent.py")
    add_logo(s)
    add_notes(s, "ADK agents use LlmAgent as the base class. The key config: model, instruction, and tools. "
              "Agents discover MCP servers via the Agent Registry using registry.get_mcp_toolset() — "
              "no hardcoded URLs. The registry resolves the server's resource name to a live connection. "
              "The Coordinator pattern uses sub_agents to delegate: it doesn't call tools directly, "
              "it routes to the specialist agent best suited for the user's intent. "
              "OTel tracing is built in — every agent call automatically generates spans.")

    # ===== SLIDE 9: DEPLOYMENT =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Deployment Architecture", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "02_deployment_architecture.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(1.5),
             "Cloud Run — MCP Servers",
             "StreamableHTTP transport, auto-scaling, IAM-secured endpoints", BLUE)
    add_card(s, Inches(7), Inches(3.1), Inches(5.5), Inches(1.5),
             "Agent Runtime — ADK Agents",
             "Vertex AI Agent Engine with built-in session management and memory", GREEN)
    txBox = add_text(s, Inches(7), Inches(5.0), Inches(5.5), Inches(1),
             "One command deploy: bash scripts/deploy_all.sh\nDeploys 3 MCP servers + coordinator agent in sequence",
             16, color=DARK)
    add_logo(s)
    add_notes(s, "Two deployment targets: Cloud Run for MCP servers (stateless HTTP services) "
              "and Agent Runtime (Vertex AI Agent Engine) for ADK agents (stateful, with session management). "
              "The deploy_all.sh script deploys everything in sequence — MCP servers first, then agents. "
              "DEMO: Run the deployment script and show the output as services come online.")

    # ===== SLIDE 10: CONSOLE — AGENT ENGINE =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Deployments on Agent Runtime", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session1_agent_engine.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_console_link(s, f"https://console.cloud.google.com/vertex-ai/agents?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Show the Agent Engine console page. Point out the deployed agent, "
              "its status, the model it's using, and the session management configuration. "
              "Highlight that Agent Runtime handles scaling, session persistence, and health checks automatically.")

    # ===== SLIDE 11: IDENTITY (SPIFFE) =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Agent Identity — SPIFFE", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "04_agent_identity_gateway.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_bullet_list(s, Inches(7), Inches(1.3), Inches(5.5), Inches(2),
                    ["Each agent gets a SPIFFE ID at deployment",
                     "Attestation policies verify agent identity at runtime",
                     "Agent Gateway enforces identity at network boundary"],
                    title="Workload Identity Federation", title_color=BLUE)
    add_table(s, Inches(7), Inches(4.0), Inches(5.5), Inches(2),
              ["Type", "Purpose"],
              [["ID-1: User", "Human accessing the app"],
               ["ID-2: Agent", "Agent's own authority"],
               ["ID-3: Delegated", "Agent on behalf of user"]],
              font_size=14)
    add_logo(s)
    add_notes(s, "SPIFFE gives every agent a cryptographic identity — not just a service account. "
              "This enables fine-grained policies: 'Agent X can call Tool Y but Agent Z cannot.' "
              "The three identity types map to real scenarios: user browsing the app (ID-1), "
              "agent autonomously running a scheduled task (ID-2), agent booking a flight for a specific user (ID-3). "
              "Agent Gateway enforces these identities at the network boundary.")

    # ===== SLIDE 12: WORKLOAD IDENTITY SETUP =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Workload Identity — Setup & Verification", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5.5), Inches(0.4),
             "Create Identity Pool & Provider", 22, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.5),
                   '# Create workload identity pool\ngcloud iam workload-identity-pools \\\n  create geap-agent-pool \\\n  --location="global" \\\n  --display-name="GEAP Agent Pool"\n\n# Create OIDC provider\ngcloud iam workload-identity-pools \\\n  providers create-oidc agent-provider \\\n  --workload-identity-pool=geap-agent-pool \\\n  --issuer-uri="https://accounts.google.com" \\\n  --location="global"',
                   font_size=11)
    add_text(s, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.5),
             "Verify token exchange: IAM & Admin > Workload Identity Federation",
             14, color=GRAY)
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(1.5),
             "Pool = Trust Boundary",
             "Each pool groups external identities from one trust domain — agents get their own pool", BLUE)
    add_card(s, Inches(7), Inches(3.1), Inches(5.5), Inches(1.5),
             "Provider = Auth Source",
             "OIDC or SAML provider that issues tokens to agents — verified at runtime by GCP", GREEN)
    add_card(s, Inches(7), Inches(4.9), Inches(5.5), Inches(1.5),
             "Token Exchange",
             "External token → STS → short-lived GCP credential — no service account keys needed", YELLOW)
    add_repo_link(s, "scripts/setup_agent_identity.sh")
    add_logo(s)
    add_notes(s, "Walk through the gcloud commands to create a workload identity pool and OIDC provider. "
              "Key concepts: the pool is a trust boundary (group external identities), "
              "the provider is the auth source (OIDC token issuer). "
              "The token exchange flow: external token goes through STS and comes back as a short-lived GCP credential. "
              "No long-lived service account keys needed — this is the security win.")

    # ===== SLIDE 13: SESSION 2 HEADER =====
    s = prs.slides.add_slide(blank)
    set_bg_gradient(s, DARK_BLUE, BLUE)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Session 2", 16, color=RGBColor(0xCC, 0xDD, 0xFF), bold=True)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
             "AI Gateway / MCP Gateway\n(continued)", 44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
             "Gateway, evaluation, observability, and optimization", 22, color=RGBColor(0xCC, 0xDD, 0xFF))
    add_bullet_list(s, Inches(0.8), Inches(4.3), Inches(10), Inches(2.5),
                    ["Agent Gateway for ingress/egress control",
                     "Three-tier evaluation pipeline",
                     "Online monitors and failure cluster analysis",
                     "Agent optimization with GEPA algorithm"],
                    font_size=20, color=RGBColor(0xEE, 0xEE, 0xFF))
    add_repo_link(s, y_pos=Inches(6.5))
    add_logo(s, RGBColor(0x88, 0x99, 0xCC))
    add_notes(s, "Session 2 builds on the deployed system from Session 1. "
              "We'll add governance (Agent Gateway controlling what agents can do), "
              "evaluation (three-tier pipeline from manual to CI/CD), "
              "observability (Cloud Trace + BigQuery analytics), "
              "and optimization (GEPA evolutionary prompt tuning). "
              "This is where the platform becomes production-ready.")

    # ===== SLIDE 14: AGENT GATEWAY =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Agent Gateway — Dual Mode", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4),
             "Ingress + Egress Governance", 22, bold=True, color=BLUE)
    add_table(s, Inches(0.8), Inches(1.6), Inches(5.8), Inches(1.5),
              ["Mode", "Gateway", "Controls"],
              [["Client-to-Agent", "geap-workshop-gateway", "Inbound user requests"],
               ["Agent-to-Anywhere", "geap-workshop-gateway-egress", "Gemini model calls, MCP tools, external APIs"]],
              font_size=13)
    txBox = add_text(s, Inches(0.8), Inches(3.4), Inches(5.8), Inches(0.8),
             "With egress gateway, ALL outbound traffic — including Gemini model calls — routes through governance (IAM Allow + SGP + Model Armor)",
             14, color=DARK)
    add_code_block(s, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.0),
                   'config = {\n  "agent_gateway_config": {\n    "client_to_agent_config": {"agent_gateway": INGRESS},\n    "agent_to_anywhere_config": {"agent_gateway": EGRESS},\n  },\n  "identity_type": "AGENT_IDENTITY",\n}',
                   font_size=11)
    add_text(s, Inches(0.8), Inches(6.5), Inches(5.8), Inches(0.4),
             "Requires Private Preview — gateway attachment via REST API PATCH",
             11, color=GRAY)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session2_agent_gateway.png"),
                   Inches(7), Inches(1.3), width=Inches(5.8))
    add_repo_link(s, "scripts/setup_agent_gateway.sh")
    add_logo(s)
    add_notes(s, "The dual-mode gateway is a key differentiator. Ingress gateway controls who can talk to your agents. "
              "Egress gateway controls what your agents can talk to — including Gemini model calls. "
              "This means ALL outbound traffic flows through governance: IAM Allow policies, "
              "Semantic Governance Policies, and Model Armor screening. "
              "Note: Egress gateway requires Private Preview access — the attachment uses a REST API PATCH. "
              "DEMO: Show the gateway configuration and the console view.")

    # ===== SLIDE 14b: AGENT TRACES =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Agent Traces", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session2_agent_traces.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Session view showing model calls, tool calls, token usage, and duration per trace",
             16, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)
    add_console_link(s, f"https://console.cloud.google.com/trace/list?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Show the Cloud Trace view of an agent interaction. "
              "Point out the trace waterfall: the root span is the user request, child spans show agent routing, "
              "model calls, and tool invocations. Highlight token usage and latency per span. "
              "This is where you debug slow responses or unexpected agent behavior.")

    # ===== SLIDE 14c: GOVERNANCE POLICIES =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Governance — Three Policy Layers", 36, bold=True)
    add_table(s, Inches(0.8), Inches(1.3), Inches(6), Inches(2.2),
              ["Layer", "Type", "Enforces"],
              [["IAM Allow", "Static (CEL rules)", "Which MCP servers an agent can call"],
               ["SGP", "Runtime (natural language)", 'Business rules — e.g. "Booking < $5,000"'],
               ["Model Armor", "Content screening", "Prompt injection, jailbreak, PII, harmful content"]],
              font_size=13)
    txBox = add_text(s, Inches(0.8), Inches(3.8), Inches(6), Inches(0.8),
             "Layered defense: IAM Allow restricts where agents connect, SGP restricts what agents do, Model Armor restricts how content is screened",
             14, color=DARK)
    add_text(s, Inches(7), Inches(1.1), Inches(5.5), Inches(0.4),
             "IAM Allow Policy (CEL)", 18, bold=True, color=BLUE)
    add_code_block(s, Inches(7), Inches(1.6), Inches(5.5), Inches(1.8),
                   'resource.type == "networkservices"\n&& resource.service in [\n  "search-mcp",\n  "booking-mcp",\n  "expense-mcp"\n]',
                   font_size=11)
    add_text(s, Inches(7), Inches(3.6), Inches(5.5), Inches(0.4),
             "Semantic Governance Policy", 18, bold=True, color=BLUE)
    add_code_block(s, Inches(7), Inches(4.1), Inches(5.5), Inches(2),
                   'name: "geap-expense-limit"\nconstraint: "Disallow expense\n  submissions exceeding $200\n  for meals, $500 for\n  entertainment"\nverdict: DENY',
                   font_size=11)
    add_repo_link(s, "scripts/setup_governance_policies.sh")
    add_logo(s)
    add_notes(s, "Three layers of governance, each catching different problems. "
              "IAM Allow is static — CEL rules that restrict which services an agent can connect to. "
              "SGP (Semantic Governance Policies) is runtime — natural language rules evaluated against live tool calls. "
              "Model Armor is content screening — catching prompt injection, jailbreaks, and harmful content. "
              "Together they form defense-in-depth: where you connect, what you do, and how content is screened. "
              "We have 5 SGP policies configured: (1) business hours restriction, (2) expense amount limits, "
              "(3) booking confirmation required, (4) anti-exfiltration guard, (5) multi-intent complexity guard. "
              "SGP verdicts are ALLOW, DENY, or ALLOW_IF_CONFIRMED (human-in-the-loop pause). "
              "Note: SGP evaluates tool calls, not raw prompts — it sees context but targets tool-call behavior.")

    # ===== SLIDE 14d: DELEGATION AUTHORIZATION =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Delegating Authorization to the Gateway", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.0), Inches(6), Inches(0.4),
             "Authz Extensions wire governance into the request path", 18, bold=True, color=BLUE)

    # Request flow diagram via table
    add_table(s, Inches(0.8), Inches(1.6), Inches(6), Inches(2.6),
              ["Extension", "Profile", "What It Sees", "Service"],
              [["IAP", "REQUEST_AUTHZ", "Headers only", "iap.googleapis.com"],
               ["Model Armor", "CONTENT_AUTHZ", "Full request + response body", "modelarmor.REGION.rep.googleapis.com"],
               ["SGP", "CONTENT_AUTHZ", "Full request + response body", "SGP engine (VPC DNS)"],
               ["Custom", "Either", "Configurable", "Your FQDN (VPC)"]],
              font_size=12)

    add_text(s, Inches(0.8), Inches(4.5), Inches(6), Inches(0.8),
             "Max 4 policies per gateway  |  REQUEST_AUTHZ runs first  |  CONTENT_AUTHZ sees bodies",
             13, color=DARK)

    # Right side: architecture flow
    add_text(s, Inches(7.2), Inches(1.0), Inches(5.5), Inches(0.4),
             "Request Flow Through Gateway", 18, bold=True, color=BLUE)
    add_code_block(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.2),
                   'Request arrives\n'
                   '  │\n'
                   '  ├─ REQUEST_AUTHZ ◄── IAP\n'
                   '  │   (identity + CEL conditions)\n'
                   '  │\n'
                   '  ├─ Route to agent/MCP\n'
                   '  │\n'
                   '  ├─ CONTENT_AUTHZ ◄── Model Armor\n'
                   '  │   (prompt/response screening)\n'
                   '  │\n'
                   '  └─ CONTENT_AUTHZ ◄── SGP\n'
                   '      (semantic business rules)',
                   font_size=11)

    # Bottom: MCP tool-level policies
    add_text(s, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.4),
             "MCP Tool-Level Policies (inline, no extension)", 16, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.5),
                   'httpRules:\n'
                   '  - to:\n'
                   '      operations:\n'
                   '        - mcp:\n'
                   '            baseProtocolMethodsOption:\n'
                   '              MATCH_BASE_PROTOCOL_METHODS\n'
                   '            methods:\n'
                   '              - name: "tools/call"\n'
                   '                params: ["search_flights"]',
                   font_size=10)

    add_text(s, Inches(7.2), Inches(5.0), Inches(5.5), Inches(0.4),
             "Setup (REST API)", 16, bold=True, color=BLUE)
    add_code_block(s, Inches(7.2), Inches(5.5), Inches(5.5), Inches(1.8),
                   '# 1. Create extension\n'
                   'POST networkservices/v1beta1/\n'
                   '  .../authzExtensions\n'
                   '  {service, failOpen, timeout}\n\n'
                   '# 2. Create policy\n'
                   'POST networksecurity/v1beta1/\n'
                   '  .../authzPolicies\n'
                   '  {target, policyProfile, action}',
                   font_size=10)

    add_logo(s)
    add_notes(s, "This slide explains the WIRING — how you connect IAP, Model Armor, and SGP to the gateway. "
              "Each governance layer runs as an authz extension, intercepting traffic at the gateway. "
              "Two profiles: REQUEST_AUTHZ sees headers only (fast, for identity checks like IAP), "
              "CONTENT_AUTHZ sees full request and response bodies (for content screening like Model Armor and SGP). "
              "Key limits: max 4 authz policies per gateway. REQUEST_AUTHZ always runs first. "
              "For Google API services (IAP, Model Armor), omit the 'authority' field and use REST API — "
              "the gcloud CLI requires an undocumented 'loadBalancingScheme' field. "
              "MCP tool-level policies can ALLOW/DENY specific tools without needing an extension. "
              "Always include baseProtocolMethodsOption: MATCH_BASE_PROTOCOL_METHODS in ALLOW rules "
              "or the MCP session will break (initialize, ping, notifications blocked). "
              "Defense-in-depth: combine REQUEST_AUTHZ (IAP for who) + CONTENT_AUTHZ (Model Armor for what content + SGP for what behavior). "
              "Docs: https://cloud.google.com/gemini-enterprise-agent-platform/govern/gateways/delegate-authorization")

    # ===== SLIDE 15: EVAL PIPELINE =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Three-Tier Evaluation Pipeline", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "03_eval_pipeline.png"),
                   Inches(3.8), Inches(1.1), height=Inches(3))
    add_card(s, Inches(0.8), Inches(4.5), Inches(3.6), Inches(2),
             "One-Time Eval",
             "Manual, on-demand evaluation with PointwiseMetric rubrics against a curated dataset", BLUE)
    add_card(s, Inches(4.8), Inches(4.5), Inches(3.6), Inches(2),
             "Online Monitors",
             "Continuous evaluation of live traffic via Cloud Trace telemetry on 10-min cycles", GREEN)
    add_card(s, Inches(8.8), Inches(4.5), Inches(3.6), Inches(2),
             "Simulated (CI/CD)",
             "Automated eval gate on PRs — score ≥ 3.0 to merge, blocks otherwise", YELLOW)
    add_repo_link(s, "src/eval/")
    add_logo(s)
    add_notes(s, "Three evaluation tiers, each serving a different purpose. "
              "One-Time Eval: manual, on-demand — you curate test cases and run PointwiseMetric rubrics. "
              "Online Monitors: continuous — evaluates live production traffic on 10-minute cycles using Cloud Trace data. "
              "Simulated (CI/CD): automated gate — blocks PRs that score below 3.0 on eval metrics. "
              "The progression is: manual validation → continuous monitoring → automated enforcement.")

    # ===== SLIDE 16: CONSOLE — EVALUATION =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Evaluation Experiments", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session2_evaluation.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_console_link(s, f"https://console.cloud.google.com/vertex-ai/agents?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Show the Evaluation Experiments page. "
              "Walk through an experiment: the eval dataset, the metrics used, and the score distribution. "
              "Point out how you can compare different agent versions side-by-side.")

    # ===== SLIDE 17: OBSERVABILITY =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Observability Stack", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "05_observability_stack.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_text(s, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
             "End-to-End Tracing", 22, bold=True, color=BLUE)
    add_text(s, Inches(7), Inches(1.9), Inches(5.5), Inches(0.6),
             "Agent Call  →  OTel Spans  →  Cloud Trace  →  BigQuery",
             18, color=BLUE, font_name=FONT)
    add_bullet_list(s, Inches(7), Inches(3.0), Inches(5.5), Inches(3.5),
                    ["Cloud Trace captures every agent + tool call",
                     "Log sink exports structured data to BigQuery",
                     "Cloud Monitoring alerts on anomalies",
                     "Failure cluster analysis groups error patterns"],
                    title="Data Pipeline", title_color=BLUE)
    add_logo(s)
    add_notes(s, "The observability pipeline: agent calls produce OTel spans automatically (built into ADK), "
              "spans are exported to Cloud Trace, a log sink copies structured data to BigQuery for analytics, "
              "and Cloud Monitoring alerts on anomalies. "
              "The failure cluster analysis is particularly useful — it groups similar errors to surface systemic issues "
              "rather than forcing you to debug individual traces.")

    # ===== SLIDE 17b: FAILURE CLUSTERS & QUALITY ALERTS =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Failure Clusters & Quality Alerts", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5.5), Inches(0.4),
             "Automated Error Analysis", 22, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.0),
                   'clusters = client.evals\n  .generate_loss_clusters(\n    src=eval_result_name\n  )\n\nfor cluster in clusters:\n  print(cluster.title)\n  print(cluster.description)\n  print(f"Samples: {cluster.sample_count}")\n  print(f"Avg score: {cluster.avg_score}")',
                   font_size=12)
    add_text(s, Inches(0.8), Inches(4.9), Inches(5.5), Inches(0.8),
             "Groups similar failures by semantic similarity — surface systemic issues instead of debugging one trace at a time",
             14, color=DARK)
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(2.0),
             "Failure Clusters",
             "Semantic grouping of eval failures into themes: tool misuse, routing errors, policy violations. "
             "Each cluster includes sample count, avg score, and representative examples.", BLUE)
    add_card(s, Inches(7), Inches(3.6), Inches(5.5), Inches(2.0),
             "Quality Alerts",
             "Cloud Monitoring alert policies fire when eval scores drop below threshold. "
             "10-minute aggregation window with custom metric: agent_eval/{metric_name}.", RED)
    add_repo_link(s, "src/eval/failure_clusters.py")
    add_logo(s)
    add_notes(s, "Failure cluster analysis is the bridge between raw eval results and actionable insights. "
              "Instead of reviewing 20+ individual failure traces, clusters group semantically similar failures: "
              "'tool misuse cluster', 'routing confusion cluster', 'policy violation cluster'. "
              "Each cluster has a title, description, sample count, and average score. "
              "Quality alerts complement this with proactive monitoring — "
              "Cloud Monitoring alert policies fire when the custom metric "
              "agent_eval/{metric_name} drops below your threshold. "
              "The 10-minute aggregation window catches quality degradation before users notice. "
              "Code: src/eval/failure_clusters.py and src/eval/quality_alerts.py")

    # ===== SLIDE 18: CONSOLE — CLOUD TRACE =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Cloud Trace", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session2_cloud_trace.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_console_link(s, f"https://console.cloud.google.com/trace/list?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Open Cloud Trace and show a real agent trace. "
              "Walk through the span hierarchy, point out model call latency, tool call duration, "
              "and token counts. This is the single most useful debugging tool for agent systems.")

    # ===== SLIDE 19: CI/CD =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "CI/CD — Simulated Eval Gate", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "06_ci_cd_flow.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_text(s, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
             "GitHub Actions Workflow", 22, bold=True, color=BLUE)
    add_text(s, Inches(7), Inches(1.9), Inches(5.5), Inches(0.6),
             "PR Opened  →  Generate Scenarios  →  Run Inference  →  Evaluate  →  Score ≥ 3.0?",
             16, color=BLUE)
    add_text(s, Inches(7), Inches(3.0), Inches(5.5), Inches(1),
             "✅  Pass — Merge allowed\n❌  Fail — PR blocked with failure report",
             18, color=DARK)
    add_logo(s)
    add_notes(s, "The CI/CD eval gate is the automated enforcement layer. "
              "When a PR is opened, GitHub Actions generates test scenarios, runs inference against the agent, "
              "evaluates responses using PointwiseMetric rubrics, and blocks the merge if score < 3.0. "
              "This prevents regressions from shipping — the same principle as unit tests but for agent quality.")

    # ===== SLIDE 20: GEPA =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Agent Optimization — GEPA Algorithm", 36, bold=True)
    add_bullet_list(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3),
                    ["Generate prompt variants using LLM mutation",
                     "Evaluate each variant against test scenarios",
                     "Select top performers based on eval scores",
                     "Evolve over multiple generations"],
                    title="Gemini Evolutionary Prompt Algorithm", title_color=BLUE)
    txBox = add_text(s, Inches(0.8), Inches(4.8), Inches(5.5), Inches(1),
             "Result: Automatically discovers better agent instructions without manual prompt engineering",
             18, color=DARK)
    add_text(s, Inches(7), Inches(1.1), Inches(5.5), Inches(0.4),
             "Configuration", 22, bold=True, color=BLUE)
    add_code_block(s, Inches(7), Inches(1.6), Inches(5.5), Inches(3.5),
                   'optimizer = AgentOptimizer(\n  agent=coordinator_agent,\n  eval_dataset=eval_data,\n  metrics=[\n    "response_quality",\n    "tool_selection",\n    "safety"\n  ],\n  generations=5,\n  population_size=4\n)',
                   font_size=13)
    add_repo_link(s, "src/optimize/run_optimize.py")
    add_logo(s)
    add_notes(s, "GEPA — Gemini Evolutionary Prompt Algorithm — automates prompt engineering. "
              "Instead of manually tuning agent instructions, GEPA generates prompt variants, "
              "evaluates each against test scenarios, selects top performers, and evolves over generations. "
              "Think of it as genetic algorithms applied to prompt optimization. "
              "The AgentOptimizer class wraps this: you provide the agent, eval data, and metrics — it finds better instructions.")

    # ===== SLIDE 20b: MULTI-MODEL ROUTER =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Multi-Model Router", 36, bold=True)
    add_text(s, Inches(0.8), Inches(1.1), Inches(5.5), Inches(0.4),
             "Complexity-Based Routing Table", 20, bold=True, color=BLUE)
    add_table(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.8),
              ["Complexity", "Model", "Score Range", "Use Case"],
              [["Low", "Flash Lite", "< 0.35", "Simple lookups"],
               ["Medium", "Gemini Flash", "0.35 – 0.64", "Multi-step queries"],
               ["High", "Claude Opus", "≥ 0.65", "Deep analysis"]],
              font_size=13)
    add_text(s, Inches(0.8), Inches(4.7), Inches(5.5), Inches(1),
             "Cost savings: 60-80% by routing simple queries to lightweight models",
             16, color=GRAY)
    add_text(s, Inches(7), Inches(1.1), Inches(5.5), Inches(0.4),
             "before_agent_callback", 20, bold=True, color=BLUE)
    add_code_block(s, Inches(7), Inches(1.6), Inches(5.5), Inches(4),
                   'async def complexity_router_callback(\n    callback_context=None, **kwargs\n):\n    user_message = ""\n    if callback_context.user_content:\n        for part in callback_context.user_content.parts:\n            if part.text:\n                user_message += part.text\n\n    result = await classify_complexity(user_message)\n    ctx = callback_context.state\n    ctx["complexity_level"] = result.level\n    ctx["complexity_score"] = result.score\n    return None',
                   font_size=12)
    add_repo_link(s, "src/router/agents.py")
    add_logo(s)
    add_notes(s, "The multi-model router uses a before_agent_callback to classify prompt complexity "
              "before the agent runs. classify_complexity() uses a lightweight Gemini call to score 0-1. "
              "Based on the score, the router delegates to lite_agent (Flash Lite), flash_agent (Flash), "
              "or opus_agent (Claude Opus via LiteLLM). This cuts costs 60-80% by sending simple queries "
              "to cheap models while reserving expensive models for complex reasoning. "
              "The callback stores the classification in session state so the router agent's instruction can read it.")

    # ===== SLIDE 20c: MEMORY BANK & USER NAMESPACES =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Memory Bank & User Namespaces", 36, bold=True)
    add_code_block(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(4.2),
                   'coordinator_agent = LlmAgent(\n    model=AGENT_MODEL,\n    name="coordinator_agent",\n    tools=[\n        get_mcp_tools(SEARCH_MCP_SERVER),\n        # Retrieves relevant memories at\n        # turn start, injects into system\n        # instruction\n        PreloadMemoryTool(),\n    ],\n    after_agent_callback=\n        save_memories_callback,\n)\n\nasync def save_memories_callback(ctx):\n    await ctx.add_session_to_memory()\n    return None',
                   font_size=12)
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(2),
             "Memory Recall",
             "PreloadMemoryTool retrieves relevant memories at turn start and injects them into the system instruction — the agent has user context before it responds", BLUE)
    add_card(s, Inches(7), Inches(3.6), Inches(5.5), Inches(2),
             "User Namespaces",
             "Memories are scoped to {user_id, app_name} — each user gets an isolated memory space with no cross-contamination", GREEN)
    add_repo_link(s, "src/agents/coordinator_agent.py")
    add_logo(s)
    add_notes(s, "Memory Bank gives agents persistent memory across sessions. Two key pieces: "
              "PreloadMemoryTool automatically retrieves relevant past interactions at turn start and "
              "injects them into the system instruction — the agent sees user context before responding. "
              "save_memories_callback runs after each turn via after_agent_callback, persisting new "
              "session events to Memory Bank for future recall. "
              "Critical: memories are scoped to {user_id, app_name} namespaces. "
              "Each user gets isolated memory — no cross-contamination between users or apps.")

    # ===== SLIDE 21: SESSION 3 HEADER =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Session 3", 16, color=YELLOW, bold=True)
    add_text(s, Inches(0.8), Inches(2), Inches(11), Inches(1.5),
             "Agent Registry", 52, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.6),
             "Agent registration, discovery, and governance", 22, color=RGBColor(0x9A, 0xA0, 0xA6))
    add_repo_link(s, y_pos=Inches(6.5))
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "Quick session on Agent Registry — the catalog for your agent fleet. "
              "If you have dozens of agents across teams, you need a way to discover, version, and govern them.")

    # ===== SLIDE 22: AGENT REGISTRY =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Agent Registry & Discovery", 36, bold=True)
    add_bullet_list(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
                    ["Register agents with metadata, toolspecs, and versioning",
                     "Discover available agents and their capabilities",
                     "Associate MCP tool specifications with agents",
                     "Govern agent lifecycle and access control"],
                    title="Registry Capabilities", title_color=BLUE)
    add_text(s, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
             "Toolspec Association", 20, bold=True, color=BLUE)
    add_code_block(s, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.5),
                   'gcloud agent-platform agent-registry \\\n  toolspecs associate \\\n  --agent=geap-coordinator \\\n  --toolspec=search-mcp-spec.yaml',
                   font_size=12)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session3_agent_registry_mcp.png"),
                   Inches(7), Inches(1.3), width=Inches(5.8))
    add_logo(s)
    add_notes(s, "Agent Registry provides four capabilities: registration (with metadata and versioning), "
              "discovery (search by capability), toolspec association (link MCP specs to agents), "
              "and lifecycle governance (access control on who can deploy or modify agents). "
              "The gcloud command shows how to associate an MCP toolspec with a registered agent. "
              "DEMO: Show the registry console and the MCP toolspec association.")

    # ===== SLIDE 23: CONSOLE — POLICIES =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Business Policies", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session3_policies.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_console_link(s, f"https://console.cloud.google.com/agent-platform/gateways?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Show the business policies page. "
              "Walk through an example policy — how it's defined, what it enforces, and how violations are handled.")

    # ===== SLIDE 24: SESSION 4 HEADER =====
    s = prs.slides.add_slide(blank)
    set_bg_gradient(s, RGBColor(0xC5, 0x22, 0x1F), RED)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Session 4", 16, color=RGBColor(0xFF, 0xCC, 0xCC), bold=True)
    add_text(s, Inches(0.8), Inches(2), Inches(11), Inches(1.5),
             "Model Security /\nModel Armor", 48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
             "Input/output screening, guardrails, and content safety", 22, color=RGBColor(0xFF, 0xCC, 0xCC))
    add_repo_link(s, y_pos=Inches(6.5))
    add_logo(s, RGBColor(0xFF, 0x99, 0x99))
    add_notes(s, "Final session — the security layer. Model Armor provides content safety for your agent system. "
              "This is critical for production: agents that talk to users need protection against prompt injection, "
              "jailbreaks, PII leakage, and harmful content generation.")

    # ===== SLIDE 25: MODEL ARMOR =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Model Armor — Content Safety Integration", 36, bold=True)
    add_image_safe(s, os.path.join(DIAGRAMS, "07_agent_armor.png"),
                   Inches(0.5), Inches(1.3), width=Inches(6))
    add_card(s, Inches(7), Inches(1.3), Inches(5.5), Inches(1.5),
             "Input Screening",
             "Prompt injection detection, jailbreak prevention, PII filtering before model call", RED)
    add_card(s, Inches(7), Inches(3.1), Inches(5.5), Inches(1.5),
             "Output Screening",
             "Harmful content filtering, hallucination flags, safety score thresholds", BLUE)
    add_card(s, Inches(7), Inches(4.9), Inches(5.5), Inches(1.5),
             "ADK Guardrail Callbacks",
             "before_model_callback and after_model_callback hooks integrated into agent lifecycle", GREEN)
    add_logo(s)
    add_notes(s, "Model Armor works at two points: input screening (before the model call) and output screening (after). "
              "Input: catches prompt injection, jailbreak attempts, PII in prompts. "
              "Output: filters harmful content, flags hallucinations, enforces safety scores. "
              "In ADK, this integrates via before_model_callback and after_model_callback — "
              "the guardrails run as part of the agent lifecycle, not as a separate service.")

    # ===== SLIDE 26: CONSOLE — MODEL ARMOR =====
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK)
    add_text(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "Console: Model Armor", 32, bold=True, color=WHITE)
    add_image_safe(s, os.path.join(SCREENSHOTS, "session4_model_armor.png"),
                   Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_console_link(s, f"https://console.cloud.google.com/security/model-armor?project={GCP_PROJECT}")
    add_logo(s, RGBColor(0x88, 0x88, 0x88))
    add_notes(s, "CONSOLE DEMO: Show the Model Armor configuration page. "
              "Walk through a template: the screening rules, the severity thresholds, "
              "and what happens when content is flagged (block vs. warn vs. log).")

    # ===== SLIDE 27: SUMMARY =====
    s = prs.slides.add_slide(blank)
    add_text(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "Platform Summary", 42, bold=True)
    cards = [
        ("Build", ["ADK agents with LlmAgent", "FastMCP tool servers", "Memory Bank + Router"], BLUE),
        ("Deploy", ["Cloud Run (MCP servers)", "Agent Runtime (agents)", "One-command deployment"], GREEN),
        ("Govern", ["SPIFFE identity", "Agent Gateway (dual)", "Agent Registry"], YELLOW),
        ("Secure", ["Model Armor templates", "Input/output screening", "Guardrail callbacks"], RED),
        ("Evaluate", ["One-time eval", "Failure clusters", "CI/CD eval gate"], BLUE),
        ("Optimize", ["GEPA algorithm", "Multi-model routing", "OTel observability"], GREEN),
    ]
    for i, (title, items, color) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(1.5 + row * 2.8)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(14)
        tf.margin_top = Pt(10)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = FONT
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = "• " + item
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY
            p2.font.name = FONT
            p2.space_before = Pt(4)
    add_logo(s)
    add_notes(s, "Recap the six pillars: Build (ADK + MCP), Deploy (Cloud Run + Agent Runtime), "
              "Govern (SPIFFE + Gateway + Registry), Secure (Model Armor), Evaluate (three-tier pipeline), "
              "Optimize (GEPA). Each pillar connects to the others — "
              "this is an integrated platform, not a collection of standalone tools.")

    # ===== SLIDE 28: RESOURCES =====
    s = prs.slides.add_slide(blank)
    set_bg_gradient(s, DARK_BLUE, BLUE)
    add_text(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
             "Resources", 16, color=RGBColor(0xCC, 0xDD, 0xFF), bold=True)
    add_text(s, Inches(0.8), Inches(1.3), Inches(11), Inches(1),
             "Get Started", 48, bold=True, color=WHITE)
    resources = [
        "Workshop repo:  github.com/jswortz/geap-tour",
        "Workshop guide:  docs/workshop_guide.md",
        "Agent Development Kit:  google.github.io/adk-docs",
        "MCP Protocol:  modelcontextprotocol.io",
        "Agent Platform:  cloud.google.com/agent-platform",
    ]
    add_bullet_list(s, Inches(0.8), Inches(2.8), Inches(10), Inches(3),
                    resources, font_size=24, color=RGBColor(0xEE, 0xEE, 0xFF))
    add_text(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
             "Thank you!", 22, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT)
    add_logo(s, RGBColor(0x88, 0x99, 0xCC))
    add_notes(s, "Share these resources with attendees. The workshop repo has all the code they ran today. "
              "The workshop guide has step-by-step instructions they can follow at their own pace. "
              "ADK docs and MCP protocol specs are the canonical references for building new agents. "
              "Thank everyone for their time and open the floor for Q&A.")

    prs.save(OUTPUT)
    print(f"Saved {len(prs.slides)} slides to {OUTPUT}")
    print(f"File size: {os.path.getsize(OUTPUT) / 1024:.0f} KB")


if __name__ == "__main__":
    build_deck()
